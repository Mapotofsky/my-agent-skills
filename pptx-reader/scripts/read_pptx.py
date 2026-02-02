import argparse
import json
import os
import sys
from typing import Optional, Tuple
from pptx import Presentation
from pptx.exc import PackageNotFoundError


def parse_bool(value: str) -> bool:
    value = value.lower()
    if value in ("true", "1", "yes", "y"):
        return True
    if value in ("false", "0", "no", "n"):
        return False
    raise ValueError("include-content 仅支持 true/false")


def normalize_range(total: int, start: Optional[int], end: Optional[int]) -> Tuple[int, int]:
    if total <= 0:
        return 0, 0
    if start is None:
        start = 1
    if end is None:
        end = total
    start = max(1, int(start))
    end = min(total, int(end))
    if start > end:
        return 0, 0
    return start - 1, end


def extract_pptx(
    file_path: str,
    slide_start: Optional[int] = None,
    slide_end: Optional[int] = None,
    include_content: bool = True,
    output_mode: str = "full"
) -> dict:
    if not os.path.isfile(file_path):
        return {
            "success": False,
            "file_path": file_path,
            "content": "",
            "slides": [],
            "statistics": {
                "slide_count": 0,
                "table_count": 0,
                "table_row_counts": [],
                "char_count": 0
            },
            "error": "文件不存在"
        }

    if not file_path.lower().endswith(".pptx"):
        return {
            "success": False,
            "file_path": file_path,
            "content": "",
            "slides": [],
            "statistics": {
                "slide_count": 0,
                "table_count": 0,
                "table_row_counts": [],
                "char_count": 0
            },
            "error": "仅支持.pptx格式"
        }

    if output_mode not in ["full", "list"]:
        return {
            "success": False,
            "file_path": file_path,
            "content": "",
            "slides": [],
            "statistics": {
                "slide_count": 0,
                "table_count": 0,
                "table_row_counts": [],
                "char_count": 0
            },
            "error": 'output_mode 仅限 ["full", "list"]'
        }

    try:
        presentation = Presentation(file_path)
    except PackageNotFoundError:
        return {
            "success": False,
            "file_path": file_path,
            "content": "",
            "slides": [],
            "statistics": {
                "slide_count": 0,
                "table_count": 0,
                "table_row_counts": [],
                "char_count": 0
            },
            "error": "无法读取pptx文件"
        }
    except Exception as exc:
        return {
            "success": False,
            "file_path": file_path,
            "content": "",
            "slides": [],
            "statistics": {
                "slide_count": 0,
                "table_count": 0,
                "table_row_counts": [],
                "char_count": 0
            },
            "error": str(exc)
        }

    slides = list(presentation.slides)
    slide_start_index, slide_end_index = normalize_range(len(slides), slide_start, slide_end)
    selected_slides = slides[slide_start_index:slide_end_index]

    content_blocks = []
    table_row_counts = []
    table_count = 0
    slides_list = []

    for index, slide in enumerate(selected_slides, start=slide_start_index + 1):
        slide_parts = []
        slide_texts = []
        slide_tables = []
        notes_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = (shape.text or "").strip()
                if text:
                    slide_parts.append(text)
                    slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    rows.append(row_cells)
                table_count += 1
                table_row_counts.append(len(rows))
                table_text = "\n".join(["\t".join(cells) for cells in rows]).strip()
                if table_text:
                    slide_parts.append(table_text)
                    slide_tables.append(table_text)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = (notes_frame.text or "").strip()
        if notes_text:
            slide_parts.append(" notes: " + notes_text)
            slide_texts.append(" notes: " + notes_text)
        if slide_parts:
            content_blocks.append("\n".join(slide_parts))
        if slide_texts:
            slide_texts = "\n".join(slide_texts).strip()
        else:
            slide_texts = ""
        slides_list.append({
            "slide_index": index,
            "texts": slide_texts,
            "tables": slide_tables
        })

    content = "\n\n".join(content_blocks).strip()
    char_count = len(content)
    if output_mode != "full" or not include_content:
        content = ""

    return {
        "success": True,
        "file_path": file_path,
        "content": content if output_mode == "full" else "",
        "slides": slides_list if output_mode == "list" else [],
        "statistics": {
            "slide_count": len(selected_slides),
            "table_count": table_count,
            "table_row_counts": table_row_counts,
            "char_count": char_count
        },
        "error": None
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("file_path")
    parser.add_argument("--slide-start", type=int, default=None)
    parser.add_argument("--slide-end", type=int, default=None)
    parser.add_argument("--include-content", type=parse_bool, default=True)
    parser.add_argument("--output-mode", choices=["full", "list"], default="full")

    args = parser.parse_args()
    result = extract_pptx(
        args.file_path,
        slide_start=args.slide_start,
        slide_end=args.slide_end,
        include_content=args.include_content,
        output_mode=args.output_mode
    )
    print(json.dumps(result, ensure_ascii=False, indent=2))

    if not result["success"]:
        sys.exit(1)


if __name__ == "__main__":
    main()
